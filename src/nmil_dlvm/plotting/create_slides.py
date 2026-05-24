import sys
import os
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[3]
if str(REPO_ROOT / "src") not in sys.path:
    sys.path.insert(0, str(REPO_ROOT / "src"))

import pandas as pd
import pingouin as pg

import torch
from nmil_dlvm.utils.variational_NN import variationalNN as variationalNN
from nmil_dlvm.utils.data_distribution_utils import (
    DATASET,
    COMPUTE_DEVICE
)
from nmil_dlvm.paths import (
    model_training_analysis_dir,
    visualization_presentations_dir,
)

import random
import math
import numpy as np
import os
import matplotlib.pyplot as plt
from matplotlib.patches import Patch
from IPython.display import display
import pickle
from numpy.linalg import norm
import ast
from matplotlib.lines import Line2D
import pdb

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu
from pptx import Presentation
from datetime import date
import pandas as pd
import os
from PIL import Image
from io import BytesIO
from torch.nn.functional import softplus

# Note: This should be done before importing any part of Pillow
Image.MAX_IMAGE_PIXELS = 500000000  # Set to a suitable value based on your image size


def create_presentation(csv_path=None, 
                        model_notes="", 
                        time_stamp=None, 
                        latent_dim=None, 
                        held_out_session_ids=None, 
                        lr=None, 
                        kld_factor=None, 
                        output_path=None,model_type = "NN"):
    if csv_path is None:
        csv_path = model_training_analysis_dir(DATASET) / "runs_data.csv"
    if output_path is None:
        output_path = visualization_presentations_dir(DATASET)

    output_path = os.fspath(output_path)
    
    # Parse the CSV and get the file_names (model_ids) based on the input parameters
    df = pd.read_csv(csv_path)

    # Convert 'time_stamp' to datetime
    df['time_stamp'] = pd.to_datetime(df['time_stamp'])

    # Filter the DataFrame based on the provided parameters
    selection = (df["held_out_session_ids"] == held_out_session_ids if held_out_session_ids is not None else True) & \
                (df["lr"] == lr if lr is not None else True)

    # If kld_factor is not provided (i.e., None), include all kld_factor values in the selection
    if kld_factor is not None:
        selection &= (df["kld_factor"] == kld_factor)

    # If latent_dim is not provided (i.e., None), include all latent_dim values in the selection
    if latent_dim is not None:
        selection &= (df["latent_dim"] == latent_dim)

    df_selected = df.loc[selection]

    # Sort the filtered df by 'time_stamp'
    df_selected = df_selected.sort_values(by='time_stamp')
    
    
    # Get the earliest time_stamp and format it to the desired string format
    earliest_time_stamp = df_selected.iloc[0]['time_stamp']
    formatted_time_stamp = earliest_time_stamp.strftime(f"DLVM-{model_type}-{DATASET}-%Y-%m-%d-(%H-%M)")

    # Use the formatted_time_stamp as the model_config
    model_config = formatted_time_stamp

    file_names = df_selected["file_name"].tolist()
    kld_factors = df_selected["kld_factor"].tolist()
    prior_loss_factors = df_selected["param_prior_factor"].tolist()
    all_runtime_args = df_selected["all_runtime_args"].tolist() #list of string dictionaries e.g. "{'held_out_session_ids': 'held_out_session_ids', 'latent_dim': 2, 'lr': 0.001, 'kld_factor': 0.1}"
    
    # if there are nans for any row, fill it with a default dictionary
    for i, arg in enumerate(all_runtime_args):
        if pd.isna(arg): # for backward compatibility
            all_runtime_args[i] = "{'n_epochs': 10000, 'l2_lambda': 0.0}"

    # import pdb; pdb.set_trace()
    all_runtime_args = [ast.literal_eval(arg) for arg in all_runtime_args]


    lrs = df_selected["lr"].tolist() # learning rates
    
    # Extract the model IDs from the file names
    model_ids = [file_name.rsplit("_", 1)[-1].split(".")[0] for file_name in file_names]

    #present models in order of increasing KLD factor
    sorted_indices = np.argsort(kld_factors)
    kld_factors = np.array(kld_factors)[sorted_indices]
    prior_loss_factors = np.array(prior_loss_factors)[sorted_indices]
    num_epochs = np.array([args["n_epochs"] for args in all_runtime_args])[sorted_indices]
    lrs = np.array(lrs)[sorted_indices]
    model_ids = np.array(model_ids)[sorted_indices]
    
    
    # Create a presentation object
    presentation = Presentation()

    # Get today's date
    today = date.today().strftime("%B %d, %Y")

    # Slide 0: Add a title slide with the model_config and date
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title = slide.shapes.title
    title.text = f"{model_config}, {DATASET} dataset,  LD = {latent_dim}, lr = {lr} Performance"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(30)

    subtitle.text = today

    # Slide 1: Add a slide with a table showing the performance metrics for each model
    metrics_table_slide = presentation.slides.add_slide(presentation.slide_layouts[5])

    # Remove any empty placeholders
    for shape in metrics_table_slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx != 0 and shape.text == "":
            sp = shape
            metrics_table_slide.shapes._spTree.remove(sp._element)
            
    table_height = Inches(3)
    table_width = Inches(7)
    # Center the table both horizontally and vertically
    left = int((presentation.slide_width - table_width) / 2)
    top = int((presentation.slide_height - table_height) / 2)

    metrics_table = metrics_table_slide.shapes.add_table(len(model_ids) + 1, 5, left, top, table_width, table_height).table
    
    # Set row headings
    metrics_table.cell(0, 0).text = "Hyperparameters"
    metrics_table.cell(0, 1).text = "Model ID"

    # After creating the table, set the width of the 'Hyperparameters' column to a new value
    # For example, increase it to 2 inches
    metrics_table.columns[0].width = Emu(3 * 914400)  # 1 inch equals 914400 EMU

    # Now, to re-center the table, we need to adjust the 'left' variable.
    # Subtract half of the added width from the 'left' position
    extra_width = Emu(3 * 914400) - metrics_table.columns[0].width
    
    # Calculate the total width of the table by summing the widths of all cells in the first row
    total_width_emu = sum(metrics_table.columns[col].width for col in range(len(metrics_table.columns)))

    # Convert EMUs to points (1 inch = 72 points)
    conversion_factor = 72.0 / 914400.0
    total_width_points = total_width_emu * conversion_factor

    # Calculate the new position of the table
    new_left = int((presentation.slide_width - total_width_points) / 2)


    # Set the new position of the table
    metrics_table.left = new_left


    # Extract the KLD Factor for each model from the original DataFrame
    for i, model_id in enumerate(model_ids, start=1):
        # Extract KLD Factor by finding the file_name that contains the model_id
        current_kld_factor = df.loc[df["file_name"].str.contains(model_id), "kld_factor"].values[0]
        model_lr = df.loc[df["file_name"].str.contains(model_id), "lr"].values[0]
        param_prior_loss_factor =  df.loc[df["file_name"].str.contains(model_id), "param_prior_factor"].values[0]
        all_runtime_args = df.loc[df["file_name"].str.contains(model_id), "all_runtime_args"].values[0] # it is a string i.e. "{'held_out_session_ids': 'held_out_session_ids', 'latent_dim': 2, 'lr': 0.001, 'kld_factor': 0.1}"
        if pd.isna(all_runtime_args): # for backward compatibility
            all_runtime_args = "{'n_epochs': 10000, 'l2_lambda': 0.0}"
        n_epoch = ast.literal_eval(all_runtime_args).get("n_epochs", 10000)
        l2_lambda = ast.literal_eval(all_runtime_args).get("l2_lambda", 0.0)

        metrics_table.cell(i, 0).text = f"KLD factor= {current_kld_factor}, LD = {latent_dim} {model_notes}, lr = {model_lr}, param_prior_factor ={param_prior_loss_factor}, num_epochs = {n_epoch},l2_lambda = {l2_lambda}"
        metrics_table.cell(i, 1).text = model_id

    # Set column headings and fill in performance metrics
    metrics = ['meu_z', 'training_logprob', 'testing_logprob']
    display_metrics_names = ['Training Meu_Z LogProbs (μ, σ)', 'Training Data Best LogProbs (μ, σ)', 'Testing Data Best LogProbs (μ, σ)']

    # Reference the /models folder infering the output path parent directory
    output_path_dirname = os.path.dirname(output_path)
    reference_path = f'./outputs/{DATASET}/' #the where the CSV and PNG files are saved


    for j, metric in enumerate(metrics, start=2):
        metrics_table.cell(0, j).text = display_metrics_names[j-2]
        for i, model_id in enumerate(model_ids, start=1):
            # Read performance metrics CSV for this model
            metrics_csv_path = os.path.join(reference_path, model_id, f"{model_id}-performance_metrics.csv")
            if os.path.exists(metrics_csv_path):
                metrics_df = pd.read_csv(metrics_csv_path)
                mu = metrics_df[f"{metric}_mean"].values[0]
                sigma = metrics_df[f"{metric}_std"].values[0]
                metrics_table.cell(i, j).text = f"({mu:.3g}, {sigma:.3g})"
            else:
                metrics_table.cell(i, j).text = "(N/A, N/A)"

        
    # Set the title of the metrics_table_slide
    
    title = metrics_table_slide.shapes.title
    title.text = "Summary Statistics"

    # Center text in each cell and reduce height to minimum possible
    for row in metrics_table.rows:
        row.height = Inches(0.15)  # Setting height as low as possible without breaking lines
        for cell in row.cells:
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(10)

    # Center text in each cell and reduce the font size
    for i in range(len(model_ids) + 1):
        for j in range(5):
            cell = metrics_table.cell(i, j)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(8)

    # Create a figure and axis with higher DPI for better resolution
    fig, ax = plt.subplots(figsize=(12, 7), dpi=200)
    
    # Set the positions of the models along the x-axis (1 to N)
    x_positions = range(1, len(model_ids) + 1)

    # Set the width of the error regions
    error_region_width = 0.15

    # Set the number of sessions being averaged for each metric
    n_sessions_training = 90 #TODO replace these numbers with the actual ones
    n_sessions_testing = 10 #TODO replace these numbers with the actual ones and add them to the plots

    # Plot mu and std for each metric as dots with transparent error regions
    metrics_labels = ['meu_z', 'training_logprob', 'testing_logprob']
    colors = ['b', 'g', 'r']
    legend_entries = []
    for i, metric in enumerate(metrics_labels):
        mu_values = []
        std_values = []
        n_sessions = n_sessions_testing if metric == 'testing_logprob' else n_sessions_training
        for model_id in model_ids:
            # Read performance metrics CSV for this model
            metrics_csv_path = os.path.join(reference_path, model_id, f"{model_id}-performance_metrics.csv")
            if os.path.exists(metrics_csv_path):
                metrics_df = pd.read_csv(metrics_csv_path)
                mu = metrics_df[f"{metric}_mean"].values[0]
                sigma = metrics_df[f"{metric}_std"].values[0]

                # Check if mu and sigma are valid numbers (not NaN)
                if pd.notna(mu) and pd.notna(sigma):
                    mu_values.append(mu)
                    std_values.append(sigma)
                else:
                    mu_values.append(None)
                    std_values.append(None)
            else:
                mu_values.append(None)
                std_values.append(None)

        # Plot mu and std as dots with transparent error regions for each metric
        line, = ax.plot([pos for pos in x_positions], mu_values, marker='o', color=colors[i], alpha=0.2)  # Adjust the alpha value for more transparency
        
        # Only compute mu - sigma and mu + sigma when both mu and sigma are not None
        lower_bound = [mu - sigma if mu is not None and sigma is not None else np.nan for mu, sigma in zip(mu_values, std_values)]
        upper_bound = [mu + sigma if mu is not None and sigma is not None else np.nan for mu, sigma in zip(mu_values, std_values)]

        ax.fill_between([pos for pos in x_positions], 
                        lower_bound, 
                        upper_bound, 
                        alpha=0.1, color=colors[i])

        legend_entries.append(line)

    # Set y lower limit to be 0 minus 5% of the y range
    y_range = ax.get_ylim()[1] - ax.get_ylim()[0]
    ax.set_ylim(-0.01, 0.3)

    # Set x-axis labels and tick positions
    ax.set_xticks([pos for pos in x_positions])
    ax.set_xticklabels([f"KLD={kld_factors[i]}" for i in range(len(kld_factors))], rotation=45, ha="right")

    # Set axis labels and title
    ax.set_xlabel('KLD Factor')
    ax.set_ylabel('Metric Value')
    ax.set_title('Metrics for Different Models')

    # Add a legend with n_sessions for each metric
    metrics_labels_with_n = [f"{label}" if label == 'testing_logprob' else f"{label}" for label in metrics_labels]
    ax.legend(legend_entries, metrics_labels_with_n)

    # Convert the figure to an image with higher DPI
    buffer = BytesIO()
    plt.tight_layout()
    plt.savefig(buffer, format='png', dpi=200)
    plt.close()
    buffer.seek(0)

    # Add a new slide for the figure
    figure_slide = presentation.slides.add_slide(presentation.slide_layouts[5])

    # Remove any empty placeholders
    for shape in figure_slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx != 0 and shape.text == "":
            sp = shape
            figure_slide.shapes._spTree.remove(sp._element)

    # Add the figure image to the slide, centering it both horizontally and vertically
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height
    img_width = Inches(10)
    img_height = Inches(5.5)
    left = int((slide_width - img_width) / 2)
    top = int((slide_height - img_height) / 1.85)
    pic = figure_slide.shapes.add_picture(buffer, left, top, width=img_width, height=img_height)

    # Add the title of the figure_slide
    title = figure_slide.shapes.title
    title.text = "Loss Trend"

    # Close the buffer
    buffer.close()


    # Add slides with images and notes for each model
    for index, model_id in enumerate(model_ids):

        # Add slides with images
        image_names = [f"{model_id}-sorted-ids-logprob_spread.png", 
                    f"{model_id}-meu_z_logprob.png", 
                    f"{model_id}-best_logprob.png", 
                    f"{model_id}-meu_z_mag.png",
                    f"{model_id}_sampled_vs_mle_params.png",
                    f"{model_id}_predicted_vs_mle_params.png",
                    f"{model_id}-all_task_data_corr_plot.png",
                    f"{model_id}-k8_primer_task_data_corr_plot.png",
                    f"{model_id}-k8_primer_vs_full_task_data_corr_plot.png",
                    f"{model_id}-k8_primer_to_full_task_data_shifts.png",
                    f"{model_id}-1obs_task_data_corr_plot.png",
                    f"{model_id}-0obs_task_data_corr_plot.png",
                    f"{model_id}_primer_vs_full_task_data_corr_plot_gif.gif",
                    f"{model_id}_primer_task_data_corr_plot_gif.gif"
                      ]
        
        # If the latent dimension is 2, add the latent-slices image
        if latent_dim == 2:
            image_names.append(f"{model_id}_primer_task_data_corr_plot_scatter_gif.gif")
            image_names.append(f"{model_id}_primer_task_data_meu_z_shifts_gif.gif")
            image_names.append(f"{model_id}-latent-slices-activation.png")
            image_names.append(f"{model_id}-latent-slices-no-activation.png")
        image_names.append(f"training_loss.png")
        
        def is_square(im):
            im_width, im_height = im.size
            aspect_ratio = im_width / im_height
            # Define a threshold for aspect ratio to consider an image square or near square
            threshold = 0.1
            
            return 1 - threshold < aspect_ratio < 1 + threshold
        for image_name in image_names:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])

            image_path = os.path.join(reference_path, model_id, image_name)

            if os.path.exists(image_path):
                # Add image to slide, centering and scaling as required
                im = Image.open(image_path)
                im_width, im_height = im.size

                if is_square(im):
                    print(image_path," is square")
                    # For square images, set height to match slide
                    slide_width = presentation.slide_width
                    slide_height = presentation.slide_height
                    img_height = slide_height
                    img_width = int(im_width * (slide_height / im_height))
                    left = int((slide_width - img_width) / 2)
                    top = 0
                else:
                    
                    slide_aspect_ratio = slide_width / slide_height

                    # Calculate the aspect ratio of the image
                    image_aspect_ratio = im_width / im_height

                    # Adjust the image dimensions to fit the slide while maintaining aspect ratio
                    if image_aspect_ratio > slide_aspect_ratio:
                        # Image is wider than the slide
                        img_width = slide_width
                        img_height = int(im_height * (slide_width / im_width))
                    else:
                        # Image is taller than or equal to the slide
                        img_height = slide_height
                        img_width = int(im_width * (slide_height / im_height))

                    # Position the image at the center of the slide
                    left = int((slide_width - img_width) / 2)
                    top = int((slide_height - img_height) / 2)

                slide.shapes.add_picture(image_path, left, top, width=img_width, height=img_height)

            else:
                # If the image doesn't exist, add a text saying "Image not found"
                text_box = slide.shapes.add_textbox(Inches(1), Inches(1), presentation.slide_width - 2*Inches(1), Inches(1))
                p = text_box.text_frame.add_paragraph()
                p.text = "Image not found"
                p.font.size = Pt(24)
                p.alignment = PP_ALIGN.CENTER

            # Add the notes textbox after the image
            notes_textbox = slide.shapes.add_textbox(0, 0, Inches(4), Inches(1))

            # Don't call add_paragraph since it's adding a new paragraph. Instead use paragraphs[0].
            notes_p = notes_textbox.text_frame.paragraphs[0]
            notes_p.text = f"KLD = {kld_factors[index]}, LD = {latent_dim}, lr ={lrs[index]} {model_notes},param_priors = {prior_loss_factors[index]}, num_epochs = {num_epochs[index]}"
            notes_p.font.size = Pt(8)  # Reduce font size to fit in textbox
            notes_p.alignment = PP_ALIGN.LEFT  # Align the text to left
            notes_textbox.text_frame.vertical_anchor = MSO_ANCHOR.TOP
            notes_textbox.text_frame.margin_top = Emu(0)  # Negative margin in EMUs
            notes_textbox.text_frame.margin_bottom = 0
            notes_textbox.text_frame.margin_left = Emu(0)  # Negative margin in EMUs
            notes_textbox.text_frame.margin_right = 0
            
    # Save the presentation
    presentation_file_name = f"{model_config}-[kld_factor={kld_factor if kld_factor is not None else 'all'},ld={latent_dim if latent_dim is not None else 'all'}, lr = {lr if lr is not None else 'all'}].pptx"
    presentation_file_path = os.path.join(output_path, presentation_file_name)
    presentation.save(presentation_file_path)
